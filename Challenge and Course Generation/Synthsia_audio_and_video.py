from pptx import Presentation
from moviepy.editor import ImageClip, concatenate_videoclips, AudioFileClip, concatenate_audioclips
from gtts import gTTS
from PIL import Image
import os

import integrated_function  

def save_slide_as_image(slide, index):
    width, height = 1280, 720
    image = Image.new('RGB', (width, height), 'white')
    
    image.save(f'slide_{index}.png')

def generate_slide_explanations(prs):
    prompts = []
    for index, slide in enumerate(prs.slides):
        prompt = f"""Provide a comprehensive audio explanation of the following slide content: {slide}."""
        prompts.append(prompt)
    
    # integrated_function.create_audio_for_prompts(prompts)

def ppt_to_video(ppt_file):
    prs = Presentation(ppt_file)
    slides = []
    audio_files = []
    
    generate_slide_explanations(prs)
    
    for index, slide in enumerate(prs.slides):
        save_slide_as_image(slide, index)
        
        audio_file = f'Shythesia_audio_mistral_{index}.mp3'
        audio_files.append(audio_file)
        
        slide_image = ImageClip(f'slide_{index}.png').set_duration(5)
        slides.append(slide_image)
    
    video = concatenate_videoclips(slides)
    
    audio_clips = [AudioFileClip(audio_file) for audio_file in audio_files]
    audio = concatenate_audioclips(audio_clips)
    video = video.set_audio(audio)
    
    video.write_videofile('Sythesia_2_video.mp4', fps=24)

    for index in range(len(prs.slides)):
        os.remove(f'slide_{index}.png')
        os.remove(f'audio_mistral_{index}.mp3')

ppt_to_video('D:\\Professional Mobility\\Video Generation\\Unit 1.pptx')
