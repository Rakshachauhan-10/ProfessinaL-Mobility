from pptx import Presentation
from moviepy.editor import ImageClip, concatenate_videoclips, AudioFileClip, concatenate_audioclips
from gtts import gTTS
from PIL import Image
import os

def save_slide_as_image(slide, index):
    width, height = 1280, 720
    image = Image.new('RGB', (width, height), 'white')
    
    image.save(f'slide_{index}.png')

def ppt_to_video(ppt_file):
    prs = Presentation(ppt_file)
    slides = []
    audio_files = []
    
    for index, slide in enumerate(prs.slides):
        save_slide_as_image(slide, index)
        
        text = f"This is slide {index + 1}"
        tts = gTTS(text)
        audio_file = f'slide_{index + 1}.mp3'
        tts.save(audio_file)
        audio_files.append(audio_file)
        
        slide_image = ImageClip(f'slide_{index}.png').set_duration(5)
        slides.append(slide_image)
    
    video = concatenate_videoclips(slides)
    
    audio_clips = [AudioFileClip(audio_file) for audio_file in audio_files]
    audio = concatenate_audioclips(audio_clips)
    video = video.set_audio(audio)
    
    video.write_videofile('Synthsia_video.mp4', fps=24)

    for index in range(len(prs.slides)):
        os.remove(f'slide_{index}.png')
        os.remove(f'slide_{index + 1}.mp3')

ppt_to_video('D:\\Professional Mobility\\Video Generation\\Unit 1.pptx')
