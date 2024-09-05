import csv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_slide(presentation, title, content):
    slide_layout = presentation.slide_layouts[1]  
    slide = presentation.slides.add_slide(slide_layout)

    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    text_frame = content_placeholder.text_frame
    text_frame.clear()  

    for line in content.split('\n'):
        p = text_frame.add_paragraph()
        p.text = line
        p.space_after = Pt(14)
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.LEFT

    return presentation

def create_presentation_from_csv(csv_file):
    prs = Presentation()
    
    with open(csv_file, newline='', encoding='utf-8') as csvfile:
        reader = csv.DictReader(csvfile)
        for row in reader:
            title = row.get('Title', '')
            content = row.get('Content', '')
            add_slide(prs, title, content)
    
    prs.save('presentation.pptx')

csv_file = 'D:\Professional Mobility\Challenge and Course Generation\course.csv'  
create_presentation_from_csv(csv_file)
