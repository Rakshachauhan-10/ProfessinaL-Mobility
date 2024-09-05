from pptx import Presentation
from pptx.util import Inches
import json

with open('course_content.json', 'r') as f:
    course_content = json.load(f)

prs = Presentation()

title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = course_content['course_title']
subtitle.text = "Course Overview and Objectives"

bullet_slide_layout = prs.slide_layouts[1]

for subtopic in course_content['subtopics']:
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = subtopic['title']
    
    tf = body_shape.text_frame
    tf.text = "Learning Objectives:"
    
    for obj in subtopic['objectives']:
        p = tf.add_paragraph()
        p.text = obj

    tf.add_paragraph()
    tf.add_paragraph().text = "Materials:"
    
    for material in subtopic['materials']:
        p = tf.add_paragraph()
        p.text = material

conclusion_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(conclusion_slide_layout)
title_shape = slide.shapes.title
body_shape = slide.placeholders[1]

title_shape.text = "Course Conclusion"
tf = body_shape.text_frame

tf.text = "Summary:"
tf.add_paragraph().text = course_content['conclusion']['summary']

tf.add_paragraph()
tf.add_paragraph().text = "Next Steps:"
for step in course_content['conclusion']['next_steps']:
    p = tf.add_paragraph()
    p.text = step

prs.save('Intermediate_Data_Science.pptx')
