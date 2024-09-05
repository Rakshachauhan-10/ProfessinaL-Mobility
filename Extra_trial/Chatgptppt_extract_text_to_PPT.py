from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

slides_content = [
    ("Portability and Advantages of Binoculars", ""),
    ("Two Facets of Portability:", "• Compactness and weight of the instrument\n• Ease of hand-holding and mounting"),
    ("Compactness and Weight", "10×50 Binoculars:\n• Common starting and adjunct binocular\n• Dimensions: ~18 cm (7 in) long and wide\n• Weight: ≤1 kg (2.2 lb), less for lightweight models"),
    ("Hand-Held Usage", "Hand-Held Use:\n• 10×50 binoculars: Easily hand-held for moderate periods\n• 15×70 or 16×70 binoculars: Hand-held for short periods\n• Larger models benefit from being mounted"),
    ("Mounting Binoculars", "Mount Options:\n• Photographic monopod or tripod with pan/tilt head\n• Suitable for binoculars up to 80 mm aperture (100 mm for lighter models)\n• Not ideal for astronomy (See Chapter 6)"),
    ("Ease of Set-Up", "Binoculars ≤100-mm Aperture:\n• Trivially easy to set up\n• Hand-held models (≤50 mm): Adjust interpupillary distance and focus\n• Ready for use immediately, no thermal equilibrium needed"),
    ("Larger Binoculars Set-Up", "Larger Binoculars:\n• Simpler to set up than telescopes\n• Typically mounted on altazimuth mounts or photographic tripods\n• Observing within 10 minutes with 6 kg binoculars on a parallelogram mount"),
    ("The Binocular Advantage", "Enhanced Detection:\n• Using two eyes improves detection of faint objects by ~1.4 times compared to one eye\n• Due to binocular summation"),
    ("Binocular Summation", "Binocular Summation:\n• Statistical Summation:\n  - Higher probability of photon detection by at least one eye\n• Likely result of multiple phenomena"),
    ("Conclusion", "Summary:\n• Binocular portability involves both physical characteristics and ease of use\n• Hand-held and mounted options available\n• Binoculars provide significant observational advantages over single-eye instruments"),
]

for title, content in slides_content:
    slide_layout = prs.slide_layouts[1]  
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    content_placeholder.text = content

prs.save('Binoculars_Presentation.pptx')
